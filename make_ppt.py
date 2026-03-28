"""Generate advisor presentation PPT."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

DARK = RGBColor(0x1A, 0x1A, 0x2E)
ACCENT = RGBColor(0x00, 0x6B, 0x9F)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF0, 0xF0, 0xF5)
GRAY_TEXT = RGBColor(0x55, 0x55, 0x55)
GREEN = RGBColor(0x2E, 0x7D, 0x32)
ORANGE = RGBColor(0xE6, 0x5C, 0x00)


def add_bg(slide, color=WHITE):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 bold=False, color=DARK, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return tf


def add_bullet_slide(slide, left, top, width, height, bullets, font_size=18,
                     color=DARK, spacing=Pt(8)):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = bullet
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Calibri"
        p.space_after = spacing
    return tf


def add_accent_bar(slide, left=0, top=1.15, width=13.333, height=0.06):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT
    shape.line.fill.background()


# ──────────────────────────────────────────────────────────────────────
# SLIDE 1: Title
# ──────────────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_bg(slide, DARK)
add_text_box(slide, 1, 1.8, 11, 1.2,
             "Tacit Collusion in Electricity Markets",
             font_size=40, bold=True, color=WHITE)
add_text_box(slide, 1, 3.2, 11, 0.8,
             "Multi-Agent PPO on a 5-Bus DC-OPF Network",
             font_size=24, color=RGBColor(0xAA, 0xCC, 0xEE))
add_text_box(slide, 1, 4.5, 11, 0.6,
             "Pranav Puranik  |  Advisor: Dr. Liu  |  March 2026",
             font_size=18, color=RGBColor(0x99, 0x99, 0xBB))
shape = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, Inches(1), Inches(4.2), Inches(3), Inches(0.04)
)
shape.fill.solid()
shape.fill.fore_color.rgb = ACCENT
shape.line.fill.background()


# ──────────────────────────────────────────────────────────────────────
# SLIDE 2: Research Question
# ──────────────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text_box(slide, 0.8, 0.4, 11, 0.8, "Research Question", font_size=32, bold=True, color=ACCENT)
add_accent_bar(slide)

add_text_box(slide, 1, 1.6, 10.5, 1.0,
             "Can independent reinforcement learning agents learn to tacitly collude\n"
             "in a realistic electricity market with network constraints?",
             font_size=22, bold=True, color=DARK)

add_bullet_slide(slide, 1, 3.0, 10.5, 4.0, [
    "Calvano et al. (2021) showed Q-learning agents collude in simple Cournot markets",
    "We extend this to: asymmetric firms, DC-OPF network, continuous actions, PPO",
    "Key metric: normalized profit gain  \u0394 = (\u03c0 \u2013 \u03c0_C) / (\u03c0_M \u2013 \u03c0_C)",
    "    \u0394 = 0: competitive outcome  |  \u0394 = 1: joint monopoly  |  \u0394 > 0: collusion",
], font_size=20, color=GRAY_TEXT)


# ──────────────────────────────────────────────────────────────────────
# SLIDE 3: Market Setup
# ──────────────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text_box(slide, 0.8, 0.4, 11, 0.8, "Market Setup: 5-Bus Network", font_size=32, bold=True, color=ACCENT)
add_accent_bar(slide)

add_bullet_slide(slide, 0.8, 1.5, 5.5, 5.5, [
    "5-bus DC network with PTDF line constraints",
    "Firm 0: 2 plants (Node 1: 150 MW, Node 2: 50 MW)",
    "Firm 1: 1 plant (Node 2: 100 MW)",
    "Quadratic cost: mc\u00b7g + \u00bd\u00b7qc\u00b7g\u00b2 per plant",
    "Linear inverse demand at each node",
    "ISO clears market via DC-OPF each step",
], font_size=18, color=GRAY_TEXT)

# Table on right
add_text_box(slide, 7, 1.5, 5.5, 0.5, "Benchmarks (computed once)", font_size=20, bold=True, color=DARK)

rows_data = [
    ["Metric", "Competitive", "Monopoly"],
    ["Avg LMP", "$23.49", "higher"],
    ["F0 profit/step", "\u03c0_C(F0)", "\u03c0_M(F0)"],
    ["F1 profit/step", "\u03c0_C(F1)", "\u03c0_M(F1)"],
    ["\u0394", "0.0", "1.0"],
]
table_shape = slide.shapes.add_table(len(rows_data), 3, Inches(7), Inches(2.2), Inches(5.3), Inches(2.5))
table = table_shape.table
for r, row in enumerate(rows_data):
    for c, val in enumerate(row):
        cell = table.cell(r, c)
        cell.text = val
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(14)
            p.font.name = "Calibri"
            p.font.color.rgb = WHITE if r == 0 else DARK
            p.alignment = PP_ALIGN.CENTER
        if r == 0:
            cell.fill.solid()
            cell.fill.fore_color.rgb = ACCENT
        else:
            cell.fill.solid()
            cell.fill.fore_color.rgb = LIGHT_GRAY if r % 2 == 0 else WHITE


# ──────────────────────────────────────────────────────────────────────
# SLIDE 4: PPO Agent Architecture
# ──────────────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text_box(slide, 0.8, 0.4, 11, 0.8, "PPO Agent Architecture", font_size=32, bold=True, color=ACCENT)
add_accent_bar(slide)

add_bullet_slide(slide, 0.8, 1.5, 5.5, 5.5, [
    "Independent PPO (IPPO): one agent per firm",
    "Actor: MLP \u2192 Gaussian mean per plant",
    "    Learnable log-std for exploration",
    "    Sigmoid \u00d7 capacity \u2192 MW dispatch",
    "Critic: MLP \u2192 scalar V(s)",
    "Shared Adam optimizer (lr = 3e-4)",
    "Architecture: 2 hidden layers, 64 units, Tanh",
], font_size=18, color=GRAY_TEXT)

add_text_box(slide, 7, 1.5, 5.5, 0.5, "Dimensions", font_size=20, bold=True, color=DARK)

dim_data = [
    ["Component", "Firm 0", "Firm 1"],
    ["Observation", "5 \u00d7 H", "5 \u00d7 H"],
    ["Action (plants)", "2", "1"],
    ["Action range", "[0,150]+[0,50]", "[0,100]"],
    ["Hidden layers", "64-64", "64-64"],
]
table_shape = slide.shapes.add_table(len(dim_data), 3, Inches(7), Inches(2.2), Inches(5.3), Inches(2.8))
table = table_shape.table
for r, row in enumerate(dim_data):
    for c, val in enumerate(row):
        cell = table.cell(r, c)
        cell.text = val
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(14)
            p.font.name = "Calibri"
            p.font.color.rgb = WHITE if r == 0 else DARK
            p.alignment = PP_ALIGN.CENTER
        if r == 0:
            cell.fill.solid()
            cell.fill.fore_color.rgb = ACCENT
        else:
            cell.fill.solid()
            cell.fill.fore_color.rgb = LIGHT_GRAY if r % 2 == 0 else WHITE


# ──────────────────────────────────────────────────────────────────────
# SLIDE 5: Training Pipeline
# ──────────────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text_box(slide, 0.8, 0.4, 11, 0.8, "Training Pipeline", font_size=32, bold=True, color=ACCENT)
add_accent_bar(slide)

pipeline_data = [
    ["Concept", "Meaning", "Value"],
    ["Step", "One env.step (1 hour of dispatch)", "~110K to converge"],
    ["Episode", "168 consecutive steps (1 week)", "168 steps"],
    ["Rollout", "Steps before one PPO update", "2,048 steps"],
    ["PPO Update", "GAE + clipped surrogate + value loss", "After each rollout"],
    ["Session", "One seed \u2192 train \u2192 converge/stop", "50 per H"],
    ["Max Timesteps", "Safety cap per session", "5,000,000"],
]
table_shape = slide.shapes.add_table(len(pipeline_data), 3, Inches(1), Inches(1.5), Inches(11), Inches(4.0))
table = table_shape.table
col_widths = [Inches(2.5), Inches(5.5), Inches(3.0)]
for i, w in enumerate(col_widths):
    table.columns[i].width = w
for r, row in enumerate(pipeline_data):
    for c, val in enumerate(row):
        cell = table.cell(r, c)
        cell.text = val
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(16)
            p.font.name = "Calibri"
            p.font.color.rgb = WHITE if r == 0 else DARK
            p.alignment = PP_ALIGN.LEFT if c == 1 else PP_ALIGN.CENTER
        if r == 0:
            cell.fill.solid()
            cell.fill.fore_color.rgb = ACCENT
        else:
            cell.fill.solid()
            cell.fill.fore_color.rgb = LIGHT_GRAY if r % 2 == 0 else WHITE

add_text_box(slide, 1, 6.0, 11, 1.0,
             "Flow:  Agents act \u2192 ISO clears (DC-OPF) \u2192 LMPs + profit \u2192 "
             "buffer stores transition \u2192 every 2048 steps: PPO update \u2192 repeat",
             font_size=16, color=GRAY_TEXT)


# ──────────────────────────────────────────────────────────────────────
# SLIDE 6: Convergence Criterion
# ──────────────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text_box(slide, 0.8, 0.4, 11, 0.8, "Convergence: \u0394-Stability Criterion", font_size=32, bold=True, color=ACCENT)
add_accent_bar(slide)

add_bullet_slide(slide, 0.8, 1.5, 11, 2.5, [
    "Calvano (Q-learning): greedy action unchanged at all states for 100,000 consecutive periods",
    "Our adaptation for continuous PPO actions:",
    "    After each PPO update, compute \u0394 for both firms from recent episode profits",
    "    If max |\u0394_new \u2013 \u0394_old| < 0.02 for all firms: increment stable_count",
    "    If stable_count \u2265 50 consecutive updates (~102K steps): declare converged",
], font_size=18, color=GRAY_TEXT)

add_text_box(slide, 0.8, 4.5, 11, 0.5, "Why not MW-based?", font_size=20, bold=True, color=DARK)
add_bullet_slide(slide, 0.8, 5.1, 11, 2.0, [
    "Continuous policies always drift slightly in raw output (\u00b10.1\u20132 MW per update)",
    "Calvano's discrete argmax either flips or doesn't \u2014 binary check. Not possible with continuous actions.",
    "\u0394-stability checks the economic outcome directly: has the profit level settled?",
], font_size=16, color=GRAY_TEXT)


# ──────────────────────────────────────────────────────────────────────
# SLIDE 7: Results H=1
# ──────────────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text_box(slide, 0.8, 0.4, 11, 0.8, "Results: H = 1  (50 sessions)", font_size=32, bold=True, color=ACCENT)
add_accent_bar(slide)

results_data = [
    ["Metric", "Firm 0 (large)", "Firm 1 (small)"],
    ["Final \u0394 (mean \u00b1 std)", "0.441 \u00b1 0.040", "1.718 \u00b1 0.092"],
    ["Interpretation", "44% of collusion gap", "Free-rides above monopoly share"],
    ["Convergence step", "~110K \u00b1 10.5K", "~110K \u00b1 10.5K"],
    ["Converged?", "All 50 sessions", "All 50 sessions"],
]
table_shape = slide.shapes.add_table(len(results_data), 3, Inches(1), Inches(1.5), Inches(11), Inches(3.0))
table = table_shape.table
col_widths = [Inches(3.5), Inches(3.75), Inches(3.75)]
for i, w in enumerate(col_widths):
    table.columns[i].width = w
for r, row in enumerate(results_data):
    for c, val in enumerate(row):
        cell = table.cell(r, c)
        cell.text = val
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(16)
            p.font.name = "Calibri"
            p.font.color.rgb = WHITE if r == 0 else DARK
            p.alignment = PP_ALIGN.CENTER
        if r == 0:
            cell.fill.solid()
            cell.fill.fore_color.rgb = ACCENT
        else:
            cell.fill.solid()
            cell.fill.fore_color.rgb = LIGHT_GRAY if r % 2 == 0 else WHITE

add_text_box(slide, 0.8, 5.0, 11, 0.5, "Key Findings", font_size=20, bold=True, color=DARK)
add_bullet_slide(slide, 0.8, 5.6, 11, 1.5, [
    "Both firms earn above competitive profits (\u0394 > 0) \u2192 tacit collusion emerges",
    "Firm 0 (larger, 200 MW) restricts output \u2192 bears the cost of restraint",
    "Firm 1 (smaller, 100 MW) free-rides on higher prices \u2192 \u0394 > 1",
    "Consistent with asymmetric Cournot theory: smaller firms benefit disproportionately",
], font_size=16, color=GRAY_TEXT)


# ──────────────────────────────────────────────────────────────────────
# SLIDE 8: Comparison with Calvano
# ──────────────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text_box(slide, 0.8, 0.4, 11, 0.8, "Comparison with Calvano et al. (2021)", font_size=32, bold=True, color=ACCENT)
add_accent_bar(slide)

comp_data = [
    ["Aspect", "Calvano", "Our Work"],
    ["Algorithm", "Q-learning (tabular)", "PPO (neural network)"],
    ["Action space", "15 discrete levels", "Continuous [0, cap] MW"],
    ["Firms", "2 symmetric", "2 asymmetric (different costs, caps)"],
    ["Network", "Single node, no constraints", "5-bus DC-OPF with line limits"],
    ["Monitoring", "Imperfect (stochastic demand)", "Imperfect (LMP history only)"],
    ["Convergence", "Argmax unchanged 100K periods", "\u0394-stable for 50 updates (~102K steps)"],
    ["Sessions", "1,000", "50"],
    ["\u0394 result", "~0.76 (symmetric)", "F0: 0.44, F1: 1.72 (asymmetric)"],
]
table_shape = slide.shapes.add_table(len(comp_data), 3, Inches(0.5), Inches(1.5), Inches(12.3), Inches(5.0))
table = table_shape.table
col_widths = [Inches(2.5), Inches(4.9), Inches(4.9)]
for i, w in enumerate(col_widths):
    table.columns[i].width = w
for r, row in enumerate(comp_data):
    for c, val in enumerate(row):
        cell = table.cell(r, c)
        cell.text = val
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(15)
            p.font.name = "Calibri"
            p.font.color.rgb = WHITE if r == 0 else DARK
            p.alignment = PP_ALIGN.LEFT if c > 0 else PP_ALIGN.CENTER
        if r == 0:
            cell.fill.solid()
            cell.fill.fore_color.rgb = ACCENT
        else:
            cell.fill.solid()
            cell.fill.fore_color.rgb = LIGHT_GRAY if r % 2 == 0 else WHITE


# ──────────────────────────────────────────────────────────────────────
# SLIDE 9: Next Steps
# ──────────────────────────────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_text_box(slide, 0.8, 0.4, 11, 0.8, "Next Steps", font_size=32, bold=True, color=ACCENT)
add_accent_bar(slide)

add_bullet_slide(slide, 0.8, 1.5, 11, 5.0, [
    "H = 2 and H = 3 jobs queued on Gilbreth (running next)",
    "    Does longer memory enable better punishment strategies?",
    "    Calvano found H=2 similar, H=3 slightly lower \u0394 (bigger state space hinders learning)",
    "",
    "Analyze limit strategies and impulse responses across sessions",
    "    Does the average strategy show Calvano-style decreasing output vs. price?",
    "",
    "Investigate the asymmetry: why \u0394(F1) >> 1",
    "    Compare with Nash-Cournot equilibrium",
    "    Does Firm 0 learn a \"Stackelberg leader\" role?",
    "",
    "Robustness: vary number of PPO epochs, learning rate, episode length",
], font_size=18, color=GRAY_TEXT)


# ──────────────────────────────────────────────────────────────────────
# Save
# ──────────────────────────────────────────────────────────────────────
out = "/Users/pranavpuranik/Documents/ppo-collusion/PPO_Collusion_Progress.pptx"
prs.save(out)
print(f"Saved → {out}")
